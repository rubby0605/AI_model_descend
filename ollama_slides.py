#!/usr/bin/env python3
"""
Generate a technical presentation on Ollama's Prefill/Decode architecture.
Output: ollama_architecture.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Color palette ──
BG_DARK    = RGBColor(0x1a, 0x1a, 0x2e)
BG_MID     = RGBColor(0x16, 0x21, 0x3e)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
ACCENT     = RGBColor(0x00, 0xD2, 0xFF)   # cyan
ACCENT2    = RGBColor(0xFF, 0x6B, 0x6B)   # coral
ACCENT3    = RGBColor(0x6B, 0xFF, 0x6B)   # green
YELLOW     = RGBColor(0xFF, 0xE6, 0x6D)
CODE_BG    = RGBColor(0x0D, 0x11, 0x17)
CODE_FG    = RGBColor(0xE6, 0xE6, 0xE6)
ORANGE     = RGBColor(0xFF, 0xA5, 0x00)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MONO_FONT = "Courier New"
BODY_FONT = "Calibri"
TITLE_FONT = "Calibri"


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text="", font_size=18,
                color=WHITE, bold=False, font_name=BODY_FONT, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    txBox.text_frame.vertical_anchor = anchor
    p = txBox.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_paragraph(text_frame, text, font_size=16, color=WHITE, bold=False,
                  font_name=BODY_FONT, align=PP_ALIGN.LEFT, space_before=Pt(4),
                  space_after=Pt(2), level=0):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    p.space_before = space_before
    p.space_after = space_after
    p.level = level
    return p


def add_code_box(slide, left, top, width, height, code_lines, font_size=11):
    """Add a dark code block with monospace text."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.fill.background()
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)

    for i, line in enumerate(code_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.font.size = Pt(font_size)
        p.font.name = MONO_FONT
        p.font.color.rgb = CODE_FG
        p.space_before = Pt(0)
        p.space_after = Pt(0)

        # Simple syntax highlighting
        stripped = line.lstrip()
        if stripped.startswith("//") or stripped.startswith("#"):
            p.font.color.rgb = RGBColor(0x6A, 0x99, 0x55)  # green comment
        elif stripped.startswith("if ") or stripped.startswith("for ") or stripped.startswith("continue"):
            p.font.color.rgb = RGBColor(0xC5, 0x86, 0xC0)  # purple keyword
        p.text = line

    return shape


def add_box(slide, left, top, width, height, text, font_size=14, fill_color=BG_MID,
            text_color=WHITE, bold=False, align=PP_ALIGN.CENTER, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = bold
    p.alignment = align
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=ACCENT, width=Pt(2)):
    connector = slide.shapes.add_connector(
        1,  # straight connector
        x1, y1, x2, y2
    )
    connector.line.color.rgb = color
    connector.line.width = width
    # Add arrowhead
    connector.end_x = x2
    connector.end_y = y2
    return connector


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=WHITE, bullet_char="\u2022"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Handle (text, highlight_color) tuples
        if isinstance(item, tuple):
            text, clr = item
        else:
            text, clr = item, color

        p.text = f"{bullet_char} {text}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = clr
        p.font.name = BODY_FONT
        p.space_before = Pt(6)
        p.space_after = Pt(4)

    return txBox


def slide_title(slide, text, subtitle=None):
    """Add a consistent slide title bar."""
    add_textbox(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                text, font_size=28, color=ACCENT, bold=True, font_name=TITLE_FONT)
    if subtitle:
        add_textbox(slide, Inches(0.6), Inches(0.9), Inches(12), Inches(0.4),
                    subtitle, font_size=14, color=LIGHT_GRAY)


def page_number(slide, num, total):
    add_textbox(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4),
                f"{num}/{total}", font_size=10, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════
# Build Presentation
# ══════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # blank

TOTAL = 12

# ── Slide 1: Title ──
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
add_textbox(sl, Inches(1), Inches(2.0), Inches(11.3), Inches(1.2),
            "Ollama 推論架構深入解析", font_size=44, color=WHITE, bold=True,
            font_name=TITLE_FONT, align=PP_ALIGN.CENTER)
add_textbox(sl, Inches(1), Inches(3.3), Inches(11.3), Inches(0.8),
            "Prefill, Decode & KV Cache", font_size=32, color=ACCENT,
            font_name=TITLE_FONT, align=PP_ALIGN.CENTER)

# Decorative line
shape = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(4.3), Inches(4.3), Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

add_textbox(sl, Inches(1), Inches(4.8), Inches(11.3), Inches(0.5),
            "Technical Deep Dive  |  假設聽眾有 ML 基礎", font_size=16, color=LIGHT_GRAY,
            align=PP_ALIGN.CENTER)
page_number(sl, 1, TOTAL)


# ── Slide 2: Transformer 推論兩階段概覽 ──
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
slide_title(sl, "Transformer 推論兩階段概覽")

# Left box — Prefill
add_box(sl, Inches(0.6), Inches(1.5), Inches(5.8), Inches(0.6),
        "Prefill（填充階段）", font_size=20, fill_color=RGBColor(0x0F, 0x3D, 0x5F),
        text_color=ACCENT, bold=True, border_color=ACCENT)

prefill_items = [
    "Prompt tokens 全部一次輸入",
    "平行計算 Attention（compute-bound）",
    "填充 KV Cache — 記住 context",
    "不生成任何 token",
]
add_bullet_list(sl, Inches(0.8), Inches(2.2), Inches(5.4), Inches(2.5),
                prefill_items, font_size=16, color=WHITE)

# Right box — Decode
add_box(sl, Inches(6.9), Inches(1.5), Inches(5.8), Inches(0.6),
        "Decode（生成階段）", font_size=20, fill_color=RGBColor(0x3D, 0x0F, 0x2F),
        text_color=ACCENT2, bold=True, border_color=ACCENT2)

decode_items = [
    "一次生成一個 token",
    "自迴歸（autoregressive）序列生成",
    "每次查詢 KV Cache（memory-bound）",
    "直到 EOS 或 max tokens",
]
add_bullet_list(sl, Inches(7.1), Inches(2.2), Inches(5.4), Inches(2.5),
                decode_items, font_size=16, color=WHITE)

# Analogy box at bottom
add_box(sl, Inches(2.5), Inches(5.3), Inches(8.3), Inches(0.9),
        "類比：Prefill = 讀懂題目（平行閱讀）  →  Decode = 逐字寫答案（序列輸出）",
        font_size=18, fill_color=RGBColor(0x2A, 0x2A, 0x40), text_color=YELLOW,
        bold=True, border_color=YELLOW)

# Arrow between boxes
add_box(sl, Inches(6.15), Inches(2.8), Inches(0.9), Inches(0.5),
        "→", font_size=28, fill_color=BG_DARK, text_color=ACCENT, bold=True)

page_number(sl, 2, TOTAL)


# ── Slide 3: Ollama 四層架構圖 ──
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
slide_title(sl, "Ollama 四層架構圖")

layers = [
    ("HTTP API Server (Go)", "routes.go", ACCENT),
    ("LLM Interface (Go)", "server.go", RGBColor(0x7B, 0xBF, 0xFF)),
    ("Runner Process (Go+cgo)", "runner.go", ACCENT3),
    ("llama.cpp (C++)", "llama.go / llama.cpp", ACCENT2),
]

y_start = Inches(1.6)
box_h = Inches(0.9)
gap = Inches(0.25)
box_w = Inches(8)
x_left = Inches(2.7)

for i, (label, file, color) in enumerate(layers):
    y = y_start + i * (box_h + gap)
    add_box(sl, x_left, y, box_w, box_h,
            label, font_size=20, fill_color=BG_MID, text_color=color,
            bold=True, border_color=color)

    # File annotation on the right
    add_textbox(sl, x_left + box_w + Inches(0.3), y, Inches(2), box_h,
                file, font_size=13, color=LIGHT_GRAY, font_name=MONO_FONT,
                anchor=MSO_ANCHOR.MIDDLE)

    # Arrow between layers
    if i < len(layers) - 1:
        arrow_y = y + box_h
        add_box(sl, Inches(6.2), arrow_y, Inches(0.9), gap,
                "▼", font_size=16, fill_color=BG_DARK, text_color=LIGHT_GRAY)

# Description
add_textbox(sl, Inches(0.6), Inches(6.2), Inches(12), Inches(0.5),
            "每一層透過函式呼叫或 cgo 橋接到下一層；最底層是 llama.cpp 的 C++ 推論引擎",
            font_size=14, color=LIGHT_GRAY)

page_number(sl, 3, TOTAL)


# ── Slide 4: Batch Processing 核心邏輯 ──
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
slide_title(sl, "Batch Processing 核心邏輯", "runner.go — processBatch()")

add_bullet_list(sl, Inches(0.6), Inches(1.3), Inches(12), Inches(1.2), [
    "每輪迴圈把所有 sequence 的待處理 tokens 累積進同一個 batch",
    "一次呼叫 Decode(batch) 送進 llama.cpp 做 forward pass",
    "Prefill 階段：整段 prompt 分批送入（受 batch size 限制）",
], font_size=16)

code = [
    "// runner.go — processBatch() 核心迴圈",
    "for i, seq := range s.seqs {",
    "    if seq == nil { continue }",
    "",
    "    for j, inp := range seq.inputs {",
    "        if batch.NumTokens() >= s.batchSize { break }",
    "",
    "        // 加入 batch：token, embedding, position, seqID, 是否要 logits",
    "        batch.Add(inp.token, inp.embed, seq.numPast,",
    "                  []int{i}, j == crossCount-1)",
    "        seq.numPast++",
    "    }",
    "    seq.inputs = seq.inputs[crossCount:]  // 移除已加入的 tokens",
    "}",
    "",
    "// 一次送入 llama.cpp 做 forward pass",
    "err := s.lc.Decode(batch)",
]
add_code_box(sl, Inches(0.6), Inches(3.0), Inches(12), Inches(4.0), code, font_size=12)

page_number(sl, 4, TOTAL)


# ── Slide 5: Prefill vs Decode 的關鍵分界 ──
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
slide_title(sl, "Prefill vs Decode 的關鍵分界", "runner.go — 判斷是否進入 sampling")

code = [
    "// runner.go:440-448",
    "for i, seq := range s.seqs {",
    "    if seq == nil { continue }",
    "",
    "    // ★ 關鍵判斷：還有未處理的 inputs → 仍在 prefill",
    "    if len(seq.inputs) != 0 {",
    "        continue  // 跳過 sampling，繼續 prefill",
    "    }",
    "",
    "    // seq.inputs 已清空 → 進入 decode → 開始 sampling",
    "    // ... sampling logic below ...",
    "}",
]
add_code_box(sl, Inches(0.6), Inches(1.5), Inches(12), Inches(3.0), code, font_size=13)

# Key insight boxes
add_box(sl, Inches(0.6), Inches(4.8), Inches(5.6), Inches(1.0),
        "seq.inputs 還有東西\n→ 還在 Prefill → 不 sample",
        font_size=16, fill_color=RGBColor(0x0F, 0x3D, 0x5F), text_color=ACCENT,
        bold=True, border_color=ACCENT)

add_box(sl, Inches(6.7), Inches(4.8), Inches(5.6), Inches(1.0),
        "seq.inputs 清空\n→ 進入 Decode → 開始 sample",
        font_size=16, fill_color=RGBColor(0x3D, 0x0F, 0x2F), text_color=ACCENT2,
        bold=True, border_color=ACCENT2)

# Summary
add_textbox(sl, Inches(0.6), Inches(6.2), Inches(12), Inches(0.6),
            "沒有顯式的 mode flag — 單純用「seq.inputs 是否為空」來切換 prefill/decode 行為",
            font_size=15, color=YELLOW, bold=True)

page_number(sl, 5, TOTAL)


# ── Slide 6: Token Generation 流程 ──
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
slide_title(sl, "Token Generation 流程（Decode Loop）", "runner.go — 自迴歸生成")

code = [
    "// runner.go:450-486 — Decode sampling",
    "token := seq.samplingCtx.Sample(s.lc, seq.iBatch)",
    "",
    "// 檢查停止條件",
    "if s.model.IsEOG(token) || seq.numPredict >= seq.predictLen {",
    "    seq.responses <- \"\"  // signal done",
    "    s.removeSequence(i)",
    "    continue",
    "}",
    "",
    "// ★ 自迴歸：把生成的 token 塞回 inputs，下一輪會被加進 batch",
    "seq.inputs = []input{{token: token}}",
    "seq.numPredict++",
]
add_code_box(sl, Inches(0.6), Inches(1.5), Inches(7.5), Inches(3.6), code, font_size=12)

# Auto-regressive loop diagram on the right
add_box(sl, Inches(8.6), Inches(1.6), Inches(4), Inches(0.7),
        "① Decode(batch)", font_size=15, fill_color=BG_MID,
        text_color=ACCENT, bold=True, border_color=ACCENT)
add_box(sl, Inches(8.6), Inches(2.6), Inches(4), Inches(0.7),
        "② Sample → token", font_size=15, fill_color=BG_MID,
        text_color=ACCENT3, bold=True, border_color=ACCENT3)
add_box(sl, Inches(8.6), Inches(3.6), Inches(4), Inches(0.7),
        "③ token → seq.inputs", font_size=15, fill_color=BG_MID,
        text_color=YELLOW, bold=True, border_color=YELLOW)
add_box(sl, Inches(8.6), Inches(4.6), Inches(4), Inches(0.7),
        "④ 回到 ① 下一輪", font_size=15, fill_color=BG_MID,
        text_color=ACCENT2, bold=True, border_color=ACCENT2)

# Loop arrow annotation
add_textbox(sl, Inches(8.8), Inches(5.5), Inches(3.8), Inches(0.8),
            "↻ 自迴歸循環，直到 EOS 或達到 max tokens",
            font_size=14, color=LIGHT_GRAY, bold=True)

page_number(sl, 6, TOTAL)


# ── Slide 7: KV Cache 架構 ──
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
slide_title(sl, "KV Cache 架構")

add_textbox(sl, Inches(0.6), Inches(1.3), Inches(12), Inches(0.6),
            "每個 Cache Slot 存儲一段對話的 Key/Value 向量，供 Decode 階段查詢",
            font_size=16, color=WHITE)

# Cache slots diagram
slot_colors = [ACCENT, RGBColor(0x7B, 0xBF, 0xFF), ACCENT3, ACCENT2]
slot_labels = [
    ("Slot 0", "對話 A — 150 tokens"),
    ("Slot 1", "對話 B — 230 tokens"),
    ("Slot 2", "空閒"),
    ("Slot 3", "對話 C — 80 tokens"),
]

for i, ((label, desc), color) in enumerate(zip(slot_labels, slot_colors)):
    y = Inches(2.2) + i * Inches(1.0)
    add_box(sl, Inches(1.5), y, Inches(2.2), Inches(0.7),
            label, font_size=16, fill_color=BG_MID, text_color=color,
            bold=True, border_color=color)

    # KV blocks visual
    block_w = Inches(0.3)
    n_blocks = [6, 9, 0, 3][i]
    for j in range(n_blocks):
        bx = Inches(4.0) + j * (block_w + Pt(4))
        shape = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx, y + Inches(0.15),
                                    block_w, Inches(0.4))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

    # Description
    desc_x = Inches(4.0) + max(n_blocks, 1) * (block_w + Pt(4)) + Inches(0.2)
    add_textbox(sl, desc_x, y, Inches(4), Inches(0.7),
                desc, font_size=13, color=LIGHT_GRAY, anchor=MSO_ANCHOR.MIDDLE)

# Prefix matching explanation
add_box(sl, Inches(0.6), Inches(6.0), Inches(12), Inches(0.9),
        "Prefix Matching：新 prompt 與已有 cache 比對前綴，若匹配可直接復用，跳過已計算的 prefill 部分",
        font_size=16, fill_color=RGBColor(0x2A, 0x2A, 0x40), text_color=YELLOW,
        bold=True, border_color=YELLOW)

page_number(sl, 7, TOTAL)


# ── Slide 8: Cache Slot 策略 ──
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
slide_title(sl, "Cache Slot 策略", "Single-user vs Multi-user")

# Single-user
add_box(sl, Inches(0.6), Inches(1.5), Inches(5.8), Inches(0.6),
        "Single-user：findLongestCacheSlot()", font_size=18,
        fill_color=RGBColor(0x0F, 0x3D, 0x5F), text_color=ACCENT, bold=True,
        border_color=ACCENT)
single_items = [
    "遍歷所有 cache slots",
    "比對新 prompt 與 slot 中已有 tokens 的前綴長度",
    "選最長 prefix match 的 slot → 最大化 cache 復用",
    "適合 CLI / 單人使用場景",
]
add_bullet_list(sl, Inches(0.8), Inches(2.2), Inches(5.4), Inches(2.5),
                single_items, font_size=15)

# Multi-user
add_box(sl, Inches(6.9), Inches(1.5), Inches(5.8), Inches(0.6),
        "Multi-user：findBestCacheSlot()", font_size=18,
        fill_color=RGBColor(0x3D, 0x0F, 0x2F), text_color=ACCENT2, bold=True,
        border_color=ACCENT2)
multi_items = [
    "同樣找最長 prefix match",
    "如果 slot 正在被其他 sequence 使用 → Cache Forking",
    "KvCacheSeqCp(src, dst, 0, longest)",
    "複製共用前綴到新 slot，各自獨立繼續",
]
add_bullet_list(sl, Inches(7.1), Inches(2.2), Inches(5.4), Inches(2.5),
                multi_items, font_size=15)

# Comparison code
code = [
    "// cache.go — findBestCacheSlot",
    "if longest > 0 && slot.inUse() {",
    "    // Fork: 複製已有的 KV cache 到新 slot",
    "    s.cache.KvCacheSeqCp(bestSlot.id, slot.id, 0, longest)",
    "    slot.numPast = longest",
    "}",
]
add_code_box(sl, Inches(0.6), Inches(5.0), Inches(12), Inches(1.8), code, font_size=12)

page_number(sl, 8, TOTAL)


# ── Slide 9: Cache Forking 圖解 ──
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
slide_title(sl, "Cache Forking 圖解")

# User A's cache
add_box(sl, Inches(1), Inches(1.6), Inches(3.5), Inches(0.7),
        "User A 的 Cache", font_size=16, fill_color=BG_MID,
        text_color=ACCENT, bold=True, border_color=ACCENT)

# Shared prefix blocks
for j in range(8):
    bx = Inches(1) + j * Inches(0.43)
    shape = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx, Inches(2.5), Inches(0.38), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT3
    shape.line.fill.background()

add_textbox(sl, Inches(1), Inches(3.1), Inches(3.5), Inches(0.4),
            "共用前綴 (prefix)", font_size=13, color=ACCENT3, bold=True,
            align=PP_ALIGN.CENTER)

# User A's unique tokens
for j in range(4):
    bx = Inches(1) + (8 + j) * Inches(0.43)
    shape = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx, Inches(2.5), Inches(0.38), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

add_textbox(sl, Inches(4.5), Inches(2.5), Inches(2), Inches(0.5),
            "← User A 獨有", font_size=12, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)

# Fork arrow
add_textbox(sl, Inches(3), Inches(3.6), Inches(2), Inches(0.8),
            "⤵ Fork", font_size=22, color=YELLOW, bold=True, align=PP_ALIGN.CENTER)

# User B's forked cache
add_box(sl, Inches(1), Inches(4.4), Inches(3.5), Inches(0.7),
        "User B 的 Cache (forked)", font_size=16, fill_color=BG_MID,
        text_color=ACCENT2, bold=True, border_color=ACCENT2)

# Forked prefix blocks
for j in range(8):
    bx = Inches(1) + j * Inches(0.43)
    shape = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx, Inches(5.3), Inches(0.38), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT3
    shape.line.fill.background()

add_textbox(sl, Inches(1), Inches(5.9), Inches(3.5), Inches(0.4),
            "複製的前綴", font_size=13, color=ACCENT3, bold=True,
            align=PP_ALIGN.CENTER)

# User B's unique tokens
for j in range(3):
    bx = Inches(1) + (8 + j) * Inches(0.43)
    shape = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx, Inches(5.3), Inches(0.38), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT2
    shape.line.fill.background()

add_textbox(sl, Inches(4.5), Inches(5.3), Inches(2), Inches(0.5),
            "← User B 獨有", font_size=12, color=ACCENT2, anchor=MSO_ANCHOR.MIDDLE)

# Code on right side
code = [
    "// cache.go:170-179",
    "s.cache.KvCacheSeqCp(",
    "    bestSlot.id,  // src: User A",
    "    slot.id,      // dst: User B",
    "    0,            // start pos",
    "    longest,      // end pos (prefix length)",
    ")",
]
add_code_box(sl, Inches(7), Inches(2.0), Inches(5.5), Inches(2.2), code, font_size=12)

add_textbox(sl, Inches(7), Inches(4.5), Inches(5.5), Inches(1.5),
            "共用系統 prompt 或相似開頭的請求時，\nCache Forking 避免重複 prefill 計算，\n大幅提升多用戶場景的吞吐量。",
            font_size=15, color=WHITE)

page_number(sl, 9, TOTAL)


# ── Slide 10: Context Window 溢出處理 ──
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
slide_title(sl, "Context Window 溢出處理", "當 numPast + inputs > numCtx")

add_bullet_list(sl, Inches(0.6), Inches(1.3), Inches(12), Inches(1.5), [
    "問題：對話太長，KV Cache 超過 context window 大小",
    "策略：保留開頭 numKeep tokens + 最近的 tokens，丟棄中間部分",
    "使用 KvCacheSeqAdd() 位移剩餘 tokens 的 position 值",
], font_size=16)

# Visual: context window diagram
# Before
add_textbox(sl, Inches(0.6), Inches(3.2), Inches(2), Inches(0.5),
            "溢出前：", font_size=15, color=WHITE, bold=True)

sections_before = [
    ("System\nPrompt", ACCENT, 2),
    ("早期對話", RGBColor(0x66, 0x66, 0x99), 4),
    ("中間對話", RGBColor(0x66, 0x66, 0x99), 4),
    ("最近對話", ACCENT2, 3),
]
x = Inches(2.5)
for label, color, n in sections_before:
    w = Inches(0.7) * n
    add_box(sl, x, Inches(3.0), w, Inches(0.7),
            label, font_size=11, fill_color=color, text_color=WHITE, bold=True)
    x += w + Pt(3)

# Arrow
add_textbox(sl, Inches(5.5), Inches(3.9), Inches(2), Inches(0.5),
            "▼  ShiftCacheSlot()", font_size=14, color=YELLOW, bold=True,
            align=PP_ALIGN.CENTER)

# After
add_textbox(sl, Inches(0.6), Inches(4.7), Inches(2), Inches(0.5),
            "溢出後：", font_size=15, color=WHITE, bold=True)

sections_after = [
    ("System\nPrompt", ACCENT, 2),
    ("丟棄 ✂", RGBColor(0x44, 0x44, 0x44), 4),
    ("最近對話", ACCENT2, 3),
]
x = Inches(2.5)
for label, color, n in sections_after:
    w = Inches(0.7) * n
    bc = ACCENT2 if "丟棄" in label else None
    add_box(sl, x, Inches(4.5), w, Inches(0.7),
            label, font_size=11, fill_color=color, text_color=WHITE, bold=True)
    x += w + Pt(3)

# Code
code = [
    "// runner.go — ShiftCacheSlot",
    "if seq.numPast + len(seq.inputs) > s.numCtx {",
    "    keep := s.numKeep         // 保留 system prompt",
    "    discard := (seq.numPast - keep) / 2  // 丟棄前半",
    "",
    "    s.lc.KvCacheSeqRm(slot, keep, keep+discard)",
    "    s.lc.KvCacheSeqAdd(slot, keep+discard, seq.numPast, -discard)",
    "    seq.numPast -= discard",
    "}",
]
add_code_box(sl, Inches(0.6), Inches(5.5), Inches(12), Inches(1.8), code, font_size=12)

page_number(sl, 10, TOTAL)


# ── Slide 11: Metrics 分離 ──
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
slide_title(sl, "Metrics 分離 — Prefill vs Decode 效能指標")

# Table header
header_y = Inches(1.6)
col_x = [Inches(0.6), Inches(4.0), Inches(7.5), Inches(10.5)]
col_w = [Inches(3.2), Inches(3.3), Inches(2.8), Inches(2.2)]
headers = ["指標名稱", "意義", "對應階段", "典型值"]

for i, (x, w, h) in enumerate(zip(col_x, col_w, headers)):
    add_box(sl, x, header_y, w, Inches(0.6), h, font_size=15,
            fill_color=RGBColor(0x2A, 0x3A, 0x5A), text_color=ACCENT, bold=True)

# Table rows
rows = [
    ["PromptEvalCount", "Prefill 處理的 token 數", "Prefill", "128-4096"],
    ["PromptEvalDuration", "Prefill 耗時", "Prefill", "50-500ms"],
    ["EvalCount", "生成的 token 數", "Decode", "1-2048"],
    ["EvalDuration", "Decode 耗時", "Decode", "1-60s"],
]

for r, row in enumerate(rows):
    ry = header_y + Inches(0.7) + r * Inches(0.65)
    row_bg = RGBColor(0x1F, 0x1F, 0x35) if r % 2 == 0 else BG_MID
    stage_color = ACCENT if row[2] == "Prefill" else ACCENT2
    for c, (x, w) in enumerate(zip(col_x, col_w)):
        tc = stage_color if c >= 2 else WHITE
        add_box(sl, x, ry, w, Inches(0.55), row[c], font_size=13,
                fill_color=row_bg, text_color=tc)

# Performance equations
add_textbox(sl, Inches(0.6), Inches(4.8), Inches(12), Inches(0.5),
            "效能計算公式", font_size=18, color=ACCENT, bold=True)

formulas = [
    ("Prefill Speed", "= PromptEvalCount / PromptEvalDuration", "tokens/sec — compute-bound, 可平行", ACCENT),
    ("Decode Speed", "= EvalCount / EvalDuration", "tokens/sec — memory-bound, 序列瓶頸", ACCENT2),
]

for i, (name, formula, desc, color) in enumerate(formulas):
    y = Inches(5.4) + i * Inches(0.9)
    add_box(sl, Inches(0.6), y, Inches(3), Inches(0.7),
            name, font_size=15, fill_color=BG_MID, text_color=color,
            bold=True, border_color=color)
    add_textbox(sl, Inches(3.8), y, Inches(3.5), Inches(0.7),
                formula, font_size=14, color=WHITE, font_name=MONO_FONT,
                anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(sl, Inches(7.5), y, Inches(5), Inches(0.7),
                desc, font_size=13, color=LIGHT_GRAY, anchor=MSO_ANCHOR.MIDDLE)

page_number(sl, 11, TOTAL)


# ── Slide 12: Summary ──
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
slide_title(sl, "Summary")

sections = [
    ("Prefill", ACCENT, [
        "批次處理所有 prompt tokens",
        "不 sample — 只填充 KV Cache",
        "Compute-bound → GPU 利用率高",
    ]),
    ("Decode", ACCENT2, [
        "自迴歸，一次生成一個 token",
        "Sample from logits → 回饋到下一輪",
        "Memory-bound → 受限於 KV Cache 讀取速度",
    ]),
    ("KV Cache", ACCENT3, [
        "Prefix matching → 最大化 cache 復用",
        "Slot 管理 + Cache Forking → 多用戶支援",
        "Context 溢出 → 智慧裁剪 (ShiftCacheSlot)",
    ]),
]

for i, (title, color, items) in enumerate(sections):
    x = Inches(0.6) + i * Inches(4.2)
    add_box(sl, x, Inches(1.5), Inches(3.8), Inches(0.7),
            title, font_size=22, fill_color=BG_MID, text_color=color,
            bold=True, border_color=color)
    add_bullet_list(sl, x + Inches(0.2), Inches(2.3), Inches(3.6), Inches(2.5),
                    items, font_size=14, color=WHITE)

# Design rationale box
add_box(sl, Inches(0.6), Inches(5.2), Inches(12), Inches(1.2),
        "為什麼分離 Prefill / Decode？\n"
        "Prefill 是 compute-bound（可平行化，榨乾 GPU）；Decode 是 memory-bound（瓶頸在 KV Cache I/O）\n"
        "分開處理讓兩個階段各自用最適合的 batching 策略，最大化硬體利用率",
        font_size=16, fill_color=RGBColor(0x2A, 0x2A, 0x40), text_color=YELLOW,
        bold=True, border_color=YELLOW, align=PP_ALIGN.LEFT)

# Footer
add_textbox(sl, Inches(1), Inches(6.8), Inches(11.3), Inches(0.4),
            "Q & A", font_size=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

page_number(sl, 12, TOTAL)


# ── Save ──
output_path = "/Users/rubylintu/code/AI_model_descend/ollama_architecture.pptx"
prs.save(output_path)
print(f"Saved {TOTAL} slides to {output_path}")
